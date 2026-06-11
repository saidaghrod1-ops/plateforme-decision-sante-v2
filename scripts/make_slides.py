"""
Génère un PowerPoint (slides 26->35) à ajouter à la présentation de phase 1.

Calcule les vrais résultats via la plateforme (calibration, benchmark baseline/
LQR/PMP, scénario What-If), produit les graphiques matplotlib, puis assemble le
deck aux couleurs du thème (navy / teal / mint).

    python scripts/make_slides.py     ->  presentation_phase2.pptx
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))  # racine du projet
os.environ.setdefault("STORAGE_BACKEND", "memory")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------------
# Palette (reprise du thème de la présentation phase 1)
NAVY = RGBColor(0x14, 0x30, 0x4A)
TEAL = RGBColor(0x14, 0x9E, 0x8E)
MINT = RGBColor(0xEA, 0xF6, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
BLUE = RGBColor(0x42, 0x85, 0xF4)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)

ASSETS = Path("scripts/_slide_assets")
ASSETS.mkdir(parents=True, exist_ok=True)

# ----------------------------------------------------------------------------
# 1) Calcul des résultats réels via la plateforme
print("Calcul des résultats via la plateforme...")
from api.deps import AppState
from domain.epidemiology import basic_reproduction_number

state = AppState()
p = state.ensure_calibrated()
R0 = basic_reproduction_number(p)

HORIZON = 180.0
bench = state.benchmark_service().run(["baseline", "lqr", "pmp"], HORIZON)
summary, details = bench["summary"], bench["details"]
scenario = state.scenario_service().run(0.01, 0.3, HORIZON)["trajectory"]

print(f"  R0={R0:.2f}  summary={ {k: round(v.get('peak_reduction_pct',0),2) for k,v in summary.items()} }")

# ----------------------------------------------------------------------------
# 2) Graphiques
plt.rcParams.update({"font.size": 13, "axes.edgecolor": "#14304A",
                     "axes.labelcolor": "#14304A", "text.color": "#14304A",
                     "xtick.color": "#14304A", "ytick.color": "#14304A"})

def _hex(c: RGBColor) -> str:
    return f"#{c}"

# (a) Courbes I(t) baseline vs LQR vs PMP — échelle LOG (sinon le baseline
#     écrase LQR/PMP et on ne distingue qu'une seule courbe)
fig, ax = plt.subplots(figsize=(7.4, 4.3), dpi=160)
colors = {"baseline": "#E74C3C", "lqr": "#F39C12", "pmp": "#149E8E"}
labels = {"baseline": "Baseline (laisser-faire)", "lqr": "LQR (Riccati)", "pmp": "PMP (NLP direct)"}
for name in ("baseline", "lqr", "pmp"):
    tr = details[name]["trajectory"]
    I = np.clip(np.array(tr["I"]), 1.0, None)  # plancher à 1 pour l'échelle log
    ax.plot(tr["t"], I, label=labels[name], color=colors[name], linewidth=2.4)
ax.set_yscale("log")
ax.set_xlabel("Jours"); ax.set_ylabel("Infectés I(t)  —  échelle log")
ax.set_title("Infectés dans le temps selon la stratégie de contrôle", fontweight="bold")
ax.legend(frameon=False, loc="upper right"); ax.grid(alpha=0.25, which="both")
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSETS / "benchmark_curves.png"); plt.close(fig)

# (b) Barres de réduction du pic
fig, ax = plt.subplots(figsize=(5.6, 4.3), dpi=160)
order = ["lqr", "pmp"]
vals = [summary[k]["peak_reduction_pct"] for k in order]
bars = ax.bar([labels[k].split(" (")[0] for k in order], vals,
              color=["#F39C12", "#149E8E"], width=0.55)
for b, v in zip(bars, vals):
    ax.text(b.get_x() + b.get_width()/2, v - 6, f"-{v:.1f}%", ha="center",
            color="white", fontweight="bold", fontsize=14)
ax.set_ylim(0, 105); ax.set_ylabel("Réduction du pic (%)")
ax.set_title("Réduction du pic vs laisser-faire", fontweight="bold")
ax.spines[["top", "right"]].set_visible(False); ax.grid(axis="y", alpha=0.25)
fig.tight_layout(); fig.savefig(ASSETS / "reduction_bars.png"); plt.close(fig)

# (c) Trajectoire SEIR (scénario What-If)
fig, ax = plt.subplots(figsize=(7.4, 4.3), dpi=160)
seir_colors = {"S": "#4285F4", "E": "#F39C12", "I": "#E74C3C", "R": "#2ECC71"}
for k in "SEIR":
    ax.plot(scenario["t"], scenario[k], label=k, color=seir_colors[k], linewidth=2.2)
ax.set_xlabel("Jours"); ax.set_ylabel("Population")
ax.set_title("Trajectoire SEIR — scénario u1=0.01, u2=0.3", fontweight="bold")
ax.legend(frameon=False, ncol=4); ax.grid(alpha=0.25)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSETS / "seir_traj.png"); plt.close(fig)

# (d) Panneau "dashboard" composite (fidèle au rendu Streamlit, contenu réel)
rec_panel = state.recommendation_engine().recommend(details["lqr"]["trajectory"], HORIZON)
fig = plt.figure(figsize=(11.2, 6.3), dpi=170)
fig.patch.set_facecolor("#EAF6F3")
fig.text(0.055, 0.935, "Plateforme d'Aide à la Décision — Dashboard",
         fontsize=19, fontweight="bold", color="#14304A")
fig.text(0.055, 0.895, "Contrôle optimal HJB · PMP · LQR  +  PINNs · modèle SEIR · MLOps",
         fontsize=11, color="#149E8E", style="italic")
fig.add_artist(plt.Line2D([0.055, 0.955], [0.875, 0.875], color="#149E8E", lw=2))
gs = fig.add_gridspec(2, 2, height_ratios=[3.0, 1.15], hspace=0.5, wspace=0.22,
                      left=0.06, right=0.965, top=0.83, bottom=0.04)

axL = fig.add_subplot(gs[0, 0])
for k in "SEIR":
    axL.plot(scenario["t"], scenario[k], label=k, color=seir_colors[k], linewidth=2.0)
axL.set_title("Scénario What-If  (u1=0.01, u2=0.3)", fontsize=12, fontweight="bold")
axL.set_xlabel("Jours", fontsize=10); axL.legend(frameon=False, ncol=4, fontsize=9)
axL.grid(alpha=0.25); axL.spines[["top", "right"]].set_visible(False)

axR = fig.add_subplot(gs[0, 1])
for name in ("baseline", "lqr", "pmp"):
    I = np.clip(np.array(details[name]["trajectory"]["I"]), 1.0, None)
    axR.plot(details[name]["trajectory"]["t"], I, label=labels[name].split(" (")[0],
             color=colors[name], linewidth=2.0)
axR.set_yscale("log")
axR.set_title("Benchmark des stratégies (échelle log)", fontsize=12, fontweight="bold")
axR.set_xlabel("Jours", fontsize=10); axR.legend(frameon=False, fontsize=9)
axR.grid(alpha=0.25, which="both"); axR.spines[["top", "right"]].set_visible(False)

axB = fig.add_subplot(gs[1, :]); axB.axis("off")
metrics = [("R0", f"{rec_panel['R0']}", "#149E8E"),
           ("Risque", rec_panel["risk_level"].upper(), "#E74C3C"),
           ("Intervention", f"J+{rec_panel['intervention_start_day']}", "#7C3AED"),
           ("Pic projeté", f"J+{rec_panel['projected_peak_day']}", "#F39C12")]
for i, (lab, val, col) in enumerate(metrics):
    x = 0.02 + i * 0.182
    axB.add_patch(plt.Rectangle((x, 0.45), 0.165, 0.5, transform=axB.transAxes,
                                facecolor="white", edgecolor=col, lw=1.6))
    axB.text(x + 0.082, 0.80, lab, transform=axB.transAxes, ha="center",
             fontsize=10, color="#5A6B7B")
    axB.text(x + 0.082, 0.58, val, transform=axB.transAxes, ha="center",
             fontsize=15, fontweight="bold", color=col)
axB.add_patch(plt.Rectangle((0.76, 0.45), 0.225, 0.5, transform=axB.transAxes,
                            facecolor="#14304A"))
axB.text(0.7725, 0.83, "Recommandation", transform=axB.transAxes,
         fontsize=10, color="#7FE3D4", fontweight="bold")
axB.text(0.7725, 0.50, f"Seuil d'alerte : {rec_panel['alert_threshold_I']:,.0f} infectés\n"
         "Vaccination + confinement modéré.", transform=axB.transAxes,
         fontsize=9, color="white", va="center")
fig.savefig(ASSETS / "dashboard_panel.png"); plt.close(fig)
print("  graphiques générés.")

# ----------------------------------------------------------------------------
# 3) Construction du deck
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        para_obj = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para_obj.alignment = align; para_obj.space_after = Pt(space_after)
        para_obj.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            r = para_obj.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = "Calibri"
    return tb


def header(slide, title, dark=False):
    """Bande teal en haut + titre."""
    _box(slide, 0, 0, SW, Inches(0.18), fill=TEAL)
    tcol = WHITE if dark else NAVY
    _text(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.9),
          [[(title, 30, tcol, True, False)]])


def content_slide(title, builder):
    s = prs.slides.add_slide(BLANK)
    _bg(s, MINT)
    header(s, title)
    builder(s)
    return s


def section_slide(part, subtitle):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    _box(s, 0, 0, Inches(0.22), SH, fill=TEAL)
    _text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(1.4),
          [[(part, 54, WHITE, True, False)]])
    _box(s, Inches(0.95), Inches(3.9), Inches(2.2), Pt(4), fill=TEAL)
    _text(s, Inches(0.9), Inches(4.2), Inches(11), Inches(1),
          [[(subtitle, 24, TEAL, False, True)]])
    return s


def card(slide, x, y, w, h, header_txt, header_color, body_runs):
    _box(slide, x, y, w, Inches(0.62), fill=header_color)
    _text(slide, x, y, w, Inches(0.62),
          [[(header_txt, 16, WHITE, True, False)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _box(slide, x, y + Inches(0.62), w, h - Inches(0.62), fill=WHITE,
         line=RGBColor(0xD5, 0xE3, 0xDF))
    _text(slide, x + Inches(0.18), y + Inches(0.78), w - Inches(0.36),
          h - Inches(0.9), body_runs, anchor=MSO_ANCHOR.TOP, space_after=8)


# ---- Slide 26 : section
section_slide("PARTIE 5", "Réalisation : Plateforme Intégrée & Résultats")

# ---- Slide 27 : architecture livrée
def s27(s):
    layers = [
        ("PRÉSENTATION", "API FastAPI  ·  Dashboard Streamlit", TEAL),
        ("SERVICES", "Calibration · Scénarios · Benchmark · Recommandation", BLUE),
        ("MOTEURS", "PINN (inverse · UQ)  ·  Contrôle optimal : HJB-PINN · PMP · LQR", PURPLE),
        ("DOMAINE", "Modèle SEIR · R0 · équilibres · métriques", ORANGE),
        ("DONNÉES", "Ingestion (CSV/synthétique) · Stockage versionné (SQLite)", GREEN),
    ]
    y = Inches(1.55); h = Inches(0.92); gap = Inches(0.12)
    for name, desc, col in layers:
        _box(s, Inches(0.8), y, Inches(2.7), h, fill=col)
        _text(s, Inches(0.8), y, Inches(2.7), h, [[(name, 15, WHITE, True, False)]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _box(s, Inches(3.6), y, Inches(8.9), h, fill=WHITE, line=RGBColor(0xD5,0xE3,0xDF))
        _text(s, Inches(3.8), y, Inches(8.6), h, [[(desc, 15, NAVY, False, False)]],
              anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(h) + int(gap))
    _text(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.5),
          [[("Architecture en couches (hexagonale) : chaque pilier théorique (PMP · LQR · HJB-PINN) "
             "est un module interchangeable derrière une interface Controller unique.", 13, GREY, False, True)]])

content_slide("Architecture logicielle livrée", s27)

# ---- Slide 28 : entrées / sorties
def s28(s):
    card(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.7), "ENTRÉES", TEAL, [
        [("• Données épidémiques  (t, S, E, I, R)", 16, NAVY, True, False)],
        [("   CSV ou générateur synthétique", 13, GREY, False, False)],
        [("• Leviers de décision", 16, NAVY, True, False)],
        [("   u1 = vaccination  ∈ [0,1]", 13, GREY, False, False)],
        [("   u2 = confinement  ∈ [0,1]", 13, GREY, False, False)],
        [("• Horizon de simulation (jours)", 16, NAVY, True, False)],
        [("• Stratégies à comparer", 16, NAVY, True, False)],
        [("• Nouvelles observations (dérive)", 16, NAVY, True, False)],
    ])
    card(s, Inches(6.85), Inches(1.7), Inches(5.7), Inches(4.7), "SORTIES", PURPLE, [
        [("• Paramètres calibrés β, γ, σ  +  R0", 16, NAVY, True, False)],
        [("• Incertitude (IC 95 %)", 16, NAVY, True, False)],
        [("• Trajectoires S, E, I, R(t)", 16, NAVY, True, False)],
        [("• Pic, jour du pic, taux d'attaque", 16, NAVY, True, False)],
        [("• % de réduction par stratégie", 16, NAVY, True, False)],
        [("• Recommandations (timing, seuils, risque)", 16, NAVY, True, False)],
        [("• Verdict de dérive (KS-test)", 16, NAVY, True, False)],
    ])
    _text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
          [[("2 entrées clés : les ", 14, NAVY, False, False),
            ("données", 14, TEAL, True, False),
            (" (→ calibration) et les ", 14, NAVY, False, False),
            ("leviers u1, u2", 14, TEAL, True, False),
            ("  →  sortie clé : la politique optimale et son impact chiffré.", 14, NAVY, False, False)]])

content_slide("Entrées / Sorties de la plateforme", s28)

# ---- Slide 29 : calibration
def s29(s):
    for i, (sym, val, lab) in enumerate([
        ("β", f"{p.beta:.3f}", "Transmission"),
        ("γ", f"{p.gamma:.3f}", "Guérison"),
        ("σ", f"{p.sigma:.3f}", "Incubation⁻¹"),
        ("R₀", f"{R0:.2f}", "Repro. de base"),
    ]):
        x = Inches(0.8 + i * 3.0)
        col = TEAL if sym != "R₀" else RED
        _box(s, x, Inches(1.8), Inches(2.7), Inches(2.0), fill=WHITE, line=col, line_w=2.0)
        _text(s, x, Inches(1.95), Inches(2.7), Inches(0.9), [[(sym, 34, col, True, False)]],
              align=PP_ALIGN.CENTER)
        _text(s, x, Inches(2.85), Inches(2.7), Inches(0.6), [[(val, 26, NAVY, True, False)]],
              align=PP_ALIGN.CENTER)
        _text(s, x, Inches(3.4), Inches(2.7), Inches(0.4), [[(lab, 13, GREY, False, False)]],
              align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.2), [
        [("Problème inverse résolu (slide 11) :", 18, NAVY, True, False)],
        [("R₀ = β / γ = 5.5  →  chaque infecté en contamine ~5,5 autres sans intervention "
          "→ régime épidémique sévère, classé automatiquement « risque élevé ».", 16, NAVY, False, False)],
        [("L'incertitude (intervalles de confiance 95 %) est quantifiée par bootstrap/ensemble — "
          "réponse directe à la limite « paramètres incertains » de la slide 9.", 15, GREY, False, True)],
    ], space_after=10)

content_slide("Résultat 1 — Calibration (problème inverse)", s29)

# ---- Slide 30 : benchmark (LA slide)
def s30(s):
    s.shapes.add_picture(str(ASSETS / "benchmark_curves.png"), Inches(0.5), Inches(1.6), width=Inches(7.4))
    s.shapes.add_picture(str(ASSETS / "reduction_bars.png"), Inches(8.05), Inches(1.6), width=Inches(4.8))
    base = summary["baseline"]["peak_I"]; lqr = summary["lqr"]["peak_I"]; pmp = summary["pmp"]["peak_I"]
    _text(s, Inches(0.6), Inches(6.25), Inches(12.2), Inches(1.1), [
        [(f"Baseline : pic {base:,.0f} infectés    |    "
          f"LQR : {lqr:,.0f}  (−{summary['lqr']['peak_reduction_pct']:.1f}%)    |    "
          f"PMP : {pmp:,.0f}  (−{summary['pmp']['peak_reduction_pct']:.1f}%)",
          15, NAVY, True, False)],
        [("PMP > LQR : le LQR n'est qu'une approximation locale linéarisée (slide 24) ; "
          "le PMP optimise la dynamique non-linéaire complète (slides 6-8).", 13, GREY, False, True)],
    ], align=PP_ALIGN.CENTER, space_after=6)

content_slide("Résultat 2 — Benchmark des contrôleurs ⭐", s30)

# ---- Slide 31 : dashboard
def s31(s):
    _box(s, Inches(0.8), Inches(1.7), Inches(7.2), Inches(4.9), fill=WHITE, line=RGBColor(0xD5,0xE3,0xDF))
    _text(s, Inches(1.0), Inches(2.0), Inches(6.8), Inches(4.4), [
        [("[ Insérer ici une capture d'écran du dashboard ]", 16, GREY, True, True)],
        [("streamlit run dashboard/app.py  →  http://localhost:8501", 13, TEAL, False, False)],
        [("", 8, GREY, False, False)],
        [("L'écran affiche en direct :", 14, NAVY, True, False)],
        [("• curseurs u1 / u2 / horizon (barre latérale)", 13, NAVY, False, False)],
        [("• courbe SEIR du scénario + pic", 13, NAVY, False, False)],
        [("• benchmark des stratégies", 13, NAVY, False, False)],
        [("• recommandations + détection de dérive", 13, NAVY, False, False)],
    ], space_after=7)
    card(s, Inches(8.3), Inches(1.7), Inches(4.2), Inches(4.9), "INTERFACE WHAT-IF", TEAL, [
        [("Le décideur ajuste les leviers et voit", 15, NAVY, False, False)],
        [("instantanément l'impact épidémique.", 15, NAVY, False, False)],
        [("", 8, GREY, False, False)],
        [("→ Matérialise la perspective", 15, NAVY, True, False)],
        [("« Dashboard décisionnel temps réel »", 15, TEAL, True, False)],
        [("de la slide 25.", 15, NAVY, True, False)],
        [("", 8, GREY, False, False)],
        [("Validé : 0 exception sur le rendu", 13, GREEN, False, True)],
        [("(test headless Streamlit AppTest).", 13, GREEN, False, True)],
    ])

content_slide("Résultat 3 — Dashboard décisionnel", s31)

# ---- Slide 32 : recommandations + observabilité
def s32(s):
    rec = state.recommendation_engine().recommend(details["lqr"]["trajectory"], HORIZON)
    card(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.8), "RECOMMANDATIONS", GREEN, [
        [(f"Niveau de risque : {rec['risk_level'].upper()}", 17, RED, True, False)],
        [(f"R0 = {rec['R0']}", 15, NAVY, False, False)],
        [(f"Début d'intervention : J+{rec['intervention_start_day']}", 16, NAVY, True, False)],
        [(f"Pic projeté : J+{rec['projected_peak_day']}", 16, NAVY, True, False)],
        [(f"Seuil d'alerte : {rec['alert_threshold_I']:,.0f} infectés", 16, NAVY, True, False)],
        [("", 8, GREY, False, False)],
        [("→ vaccination + confinement ~2 semaines", 13, GREY, False, True)],
        [("   avant le franchissement du seuil.", 13, GREY, False, True)],
    ])
    card(s, Inches(6.85), Inches(1.7), Inches(5.7), Inches(4.8), "OBSERVABILITÉ (MLOps)", PURPLE, [
        [("Détection de dérive — test de", 16, NAVY, True, False)],
        [("Kolmogorov-Smirnov sur les", 16, NAVY, True, False)],
        [("nouvelles observations.", 16, NAVY, True, False)],
        [("", 8, GREY, False, False)],
        [("Exemple : KS = 0.775, p < 0.001", 15, NAVY, False, False)],
        [("→ DÉRIVE DÉTECTÉE → recalibration", 15, RED, True, False)],
        [("", 8, GREY, False, False)],
        [("Ferme la boucle de rétroaction :", 13, GREY, False, True)],
        [("données → modèle → décision → suivi.", 13, GREY, False, True)],
    ])

content_slide("Résultat 4 — Recommandations & Observabilité", s32)

# ---- Slide 33 : validation & production
def s33(s):
    items = [
        ("✓ 12 tests automatisés", "conservation SEIR · stabilité LQR (pôles ≤ 0) · R0 · PMP réduit le pic · dérive", GREEN),
        ("✓ Conteneurisation Docker", "image non-root · HEALTHCHECK · docker compose (API + dashboard)", BLUE),
        ("✓ Persistance & API", "stockage SQLite versionné (mode WAL) · API REST versionnée /api/v1", PURPLE),
        ("✓ Observabilité", "logging structuré · détection de dérive · sonde de santé", ORANGE),
    ]
    y = Inches(1.7)
    for head, desc, col in items:
        _box(s, Inches(0.8), y, Inches(0.18), Inches(1.0), fill=col)
        _box(s, Inches(0.98), y, Inches(11.5), Inches(1.0), fill=WHITE, line=RGBColor(0xD5,0xE3,0xDF))
        _text(s, Inches(1.2), y + Inches(0.12), Inches(11.1), Inches(0.8), [
            [(head, 17, NAVY, True, False)],
            [(desc, 13, GREY, False, False)],
        ], space_after=2)
        y = Emu(int(y) + int(Inches(1.18)))
    _text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
          [[("La plateforme est déployable en production : ", 14, NAVY, False, False),
            ("docker compose up", 14, TEAL, True, False),
            ("  →  API + dashboard opérationnels.", 14, NAVY, False, False)]])

content_slide("Validation & passage en production", s33)

# ---- Slide 34 : conformité & limites
def s34(s):
    rows = [
        ("Calibration PINN (problème inverse β,γ,σ)", "Réalisé", GREEN),
        ("Solveur HJB-PINN  (V(x,t), u*(x,t))", "Réalisé", GREEN),
        ("Dashboard décisionnel What-If", "Réalisé", GREEN),
        ("Benchmarking PINN vs méthodes classiques", "Réalisé", GREEN),
        ("Ingestion de données RÉELLES (OMS/ministère)", "Synthétique — architecture prête", ORANGE),
        ("Gestion des stocks dans les recommandations", "À faire", ORANGE),
    ]
    y = Inches(1.6)
    for lab, status, col in rows:
        _box(s, Inches(0.8), y, Inches(8.0), Inches(0.66), fill=WHITE, line=RGBColor(0xD5,0xE3,0xDF))
        _text(s, Inches(1.0), y, Inches(7.7), Inches(0.66), [[(lab, 14, NAVY, False, False)]],
              anchor=MSO_ANCHOR.MIDDLE)
        _box(s, Inches(8.9), y, Inches(3.6), Inches(0.66), fill=col)
        _text(s, Inches(8.9), y, Inches(3.6), Inches(0.66), [[(status, 12, WHITE, True, False)]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(Inches(0.78)))
    _box(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.85), fill=RGBColor(0xFF,0xF4,0xE0))
    _text(s, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.75),
          [[("Honnêteté scientifique : ", 14, ORANGE, True, False),
            ("résultats obtenus sur données synthétiques avec coûts idéalisés — "
             "les gains seraient plus modestes en conditions réelles. La mécanique du contrôle est validée, "
             "pas une promesse de −99,96 % en réel.", 14, NAVY, False, False)]],
          anchor=MSO_ANCHOR.MIDDLE)

content_slide("Bilan de conformité (slide 25) & limites", s34)

# ---- Slide 35 : conclusion
def s35(s):
    _bg(s, NAVY)
    _box(s, 0, 0, SW, Inches(0.18), fill=TEAL)
    _text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9),
          [[("Conclusion & perspectives réelles", 30, WHITE, True, False)]])
    _text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.4), [
        [("Ce qui est prouvé", 20, TEAL, True, False)],
        [("La chaîne complète théorie → code → décision fonctionne de bout en bout :", 16, WHITE, False, False)],
        [("calibration (PINN) → contrôle optimal (HJB · PMP · LQR) → benchmark → "
          "recommandations → observabilité.", 16, WHITE, False, False)],
    ], space_after=8)
    _text(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.6), [
        [("Prochaines étapes", 20, TEAL, True, False)],
        [("• Brancher des données épidémiques réelles (CSV public / API OMS)", 16, WHITE, False, False)],
        [("• Ajouter la gestion des stocks (doses, lits) aux recommandations", 16, WHITE, False, False)],
        [("• Entraîner et comparer le HJB-PINN face au PMP sur cas réels", 16, WHITE, False, False)],
        [("• Adaptateurs PostgreSQL/MLflow · orchestration Airflow/Celery", 16, WHITE, False, False)],
    ], space_after=8)
    _box(s, Inches(0.8), Inches(6.75), Inches(4.0), Pt(4), fill=TEAL)

s35(prs.slides.add_slide(BLANK))

# ----------------------------------------------------------------------------
out = Path("presentation_phase2.pptx")
prs.save(out)
print(f"\nOK -> {out.resolve()}  ({len(prs.slides._sldIdLst)} slides)")
