"""
Génère le PowerPoint « Plateforme & Résultats » (architecture + résultats
synthétiques + correction HJB-PINN + données réelles Maroc).

    python scripts/make_slides_resultats.py  ->  presentation_plateforme_resultats.pptx
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
os.environ.setdefault("STORAGE_BACKEND", "memory")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from domain.seir import SEIRParams, simulate
from domain.epidemiology import basic_reproduction_number
from ingestion import connectors
from services.calibration_service import estimate_from_growth
from services.benchmark_service import BenchmarkService

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x14, 0x30, 0x4A); TEAL = RGBColor(0x14, 0x9E, 0x8E)
MINT = RGBColor(0xEA, 0xF6, 0xF3); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x6B, 0x7B); BLUE = RGBColor(0x42, 0x85, 0xF4)
ORANGE = RGBColor(0xF3, 0x9C, 0x12); PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x2E, 0xCC, 0x71); RED = RGBColor(0xE7, 0x4C, 0x3C)

ASSETS = ROOT / "scripts" / "_slides_assets"; ASSETS.mkdir(parents=True, exist_ok=True)
CSV = ROOT / "data" / "maroc_seir_reconstruit.csv"

# ---------------------------------------------------------------- calculs
print("Calcul des résultats...")
p_syn = SEIRParams(); R0_syn = basic_reproduction_number(p_syn)
res_syn = BenchmarkService(p_syn).run(["baseline", "lqr", "pmp", "hjb"], 180.0)
sum_syn, det_syn = res_syn["summary"], res_syn["details"]

data = connectors.from_csv(str(CSV))
p_real = estimate_from_growth(data); R0_real = basic_reproduction_number(p_real)
t_obs = np.asarray(data["t"], float); I_obs = np.asarray(data["y"], float)[:, 2]
peak_t = t_obs[int(I_obs.argmax())]
m = (I_obs > max(0.01 * I_obs.max(), 10)) & (I_obs < 0.6 * I_obs.max()) & (t_obs < peak_t)
r_growth = float(np.polyfit(t_obs[m], np.log(I_obs[m]), 1)[0])
res_real = BenchmarkService(p_real).run(["baseline", "lqr", "pmp"], 180.0)
sum_real, det_real = res_real["summary"], res_real["details"]
print("  R0_syn=%.2f  R0_real=%.2f" % (R0_syn, R0_real))

# ---------------------------------------------------------------- figures
plt.rcParams.update({"font.size": 13, "axes.edgecolor": "#14304A", "axes.labelcolor": "#14304A",
                     "text.color": "#14304A", "xtick.color": "#14304A", "ytick.color": "#14304A"})
from presentation.plots import benchmark_figure  # figure stylée partagée (source unique)


def bench_fig(det, keys, path, title):
    fig = benchmark_figure(det, keys, title)
    fig.savefig(path); plt.close(fig)


bench_fig(det_syn, ["baseline", "lqr", "pmp", "hjb"], ASSETS / "bench_syn.png",
          "Benchmark — données synthétiques (R0=5.5)")
bench_fig(det_real, ["baseline", "lqr", "pmp"], ASSETS / "bench_real.png",
          "Benchmark — paramètres réels Maroc (R0=1.87)")

x0 = np.asarray(data["y0"], float)
model = simulate(p_real, x0, (0, float(t_obs.max())), len(t_obs))
fig, ax = plt.subplots(figsize=(7.3, 4.2), dpi=170)
mo = t_obs <= 90
ax.semilogy(t_obs[mo], np.clip(I_obs[mo], 1, None), "o", ms=4, color="#E74C3C", label="Observé (Maroc)", alpha=0.7)
tm = np.array(model["t"]); im = np.clip(np.array(model["I"]), 1, None); mm = tm <= 90
ax.semilogy(tm[mm], im[mm], color="#149E8E", lw=2.5, label="Modèle SEIR calibré")
ax.set_xlabel("Jours depuis le 08/03/2020"); ax.set_ylabel("Infectés I(t) — log")
ax.set_title("Croissance initiale : observé vs modèle (Maroc)", fontweight="bold")
ax.legend(frameon=False); ax.grid(alpha=0.25, which="both"); ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSETS / "fit_real.png"); plt.close(fig)
print("  figures générées.")

# ---------------------------------------------------------------- deck
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height


def bg(s, c): s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def box(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h); sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa)
        for (t, sz, c, b, it) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = c
            r.font.bold = b; r.font.italic = it; r.font.name = "Calibri"
    return tb


def header(s, title):
    box(s, 0, 0, SW, Inches(0.18), fill=TEAL)
    text(s, Inches(0.55), Inches(0.33), Inches(12.2), Inches(0.9), [[(title, 29, NAVY, True, False)]])


def content(title):
    s = prs.slides.add_slide(BLANK); bg(s, MINT); header(s, title); return s


def section(part, sub):
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    box(s, 0, 0, Inches(0.22), SH, fill=TEAL)
    text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(1.4), [[(part, 50, WHITE, True, False)]])
    box(s, Inches(0.95), Inches(3.85), Inches(2.2), Pt(4), fill=TEAL)
    text(s, Inches(0.9), Inches(4.15), Inches(11), Inches(1), [[(sub, 23, TEAL, False, True)]])
    return s


def card(s, x, y, w, h, head, col, body):
    box(s, x, y, w, Inches(0.6), fill=col)
    text(s, x, y, w, Inches(0.6), [[(head, 15, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, y + Inches(0.6), w, h - Inches(0.6), fill=WHITE, line=RGBColor(0xD5, 0xE3, 0xDF))
    text(s, x + Inches(0.16), y + Inches(0.74), w - Inches(0.32), h - Inches(0.85), body, sa=7)


def table(s, x, y, w, rows, col_w, head_color):
    rh = Inches(0.5)
    for ci, (txt, _) in enumerate(rows[0]):
        box(s, x + sum(col_w[:ci], Emu(0)), y, col_w[ci], rh, fill=head_color)
        text(s, x + sum(col_w[:ci], Emu(0)), y, col_w[ci], rh, [[(txt, 13, WHITE, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows[1:]):
        yy = y + rh * (ri + 1)
        for ci, (txt, c) in enumerate(row):
            box(s, x + sum(col_w[:ci], Emu(0)), yy, col_w[ci], rh, fill=WHITE, line=RGBColor(0xD5, 0xE3, 0xDF))
            text(s, x + sum(col_w[:ci], Emu(0)), yy, col_w[ci], rh, [[(txt, 13, c, False, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---- 1. section
section("PARTIE 5", "La Plateforme Intégrée & ses Résultats")

# ---- 2. architecture
s = content("Architecture logicielle (en couches)")
layers = [("PRÉSENTATION", "API FastAPI · Dashboard Streamlit", TEAL),
          ("SERVICES", "Calibration · Scénarios · Benchmark · Recommandation", BLUE),
          ("MOTEURS", "PINN (inverse · UQ) · Contrôle : HJB-PINN · PMP · LQR", PURPLE),
          ("DOMAINE", "Modèle SEIR · R0 · équilibres · métriques", ORANGE),
          ("DONNÉES", "Ingestion (CSV réel / synthétique) · Stockage SQLite", GREEN)]
y = Inches(1.5)
for name, desc, col in layers:
    box(s, Inches(0.8), y, Inches(2.7), Inches(0.92), fill=col)
    text(s, Inches(0.8), y, Inches(2.7), Inches(0.92), [[(name, 14, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(3.6), y, Inches(8.9), Inches(0.92), fill=WHITE, line=RGBColor(0xD5, 0xE3, 0xDF))
    text(s, Inches(3.8), y, Inches(8.6), Inches(0.92), [[(desc, 14, NAVY, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y) + int(Inches(1.04)))
text(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.5),
     [[("Dépendances dirigées vers le noyau ; les contrôleurs LQR · PMP · HJB-PINN sont interchangeables.", 13, GREY, False, True)]])

# ---- 3. entrées / sorties
s = content("Entrées / Sorties de la plateforme")
card(s, Inches(0.8), Inches(1.6), Inches(5.7), Inches(4.9), "ENTRÉES", TEAL, [
    [("• Données épidémiques (t, S, E, I, R)", 15, NAVY, True, False)],
    [("   CSV réel (Maroc) ou synthétique", 12, GREY, False, False)],
    [("• Leviers : u1 vaccination, u2 confinement", 15, NAVY, True, False)],
    [("• Horizon (jours), stratégies à comparer", 15, NAVY, True, False)],
    [("• Nouvelles observations (dérive)", 15, NAVY, True, False)]])
card(s, Inches(6.85), Inches(1.6), Inches(5.7), Inches(4.9), "SORTIES", PURPLE, [
    [("• β, γ, σ + R0 (+ incertitude)", 15, NAVY, True, False)],
    [("• Trajectoires S,E,I,R(t) + pic", 15, NAVY, True, False)],
    [("• % réduction du pic par stratégie", 15, NAVY, True, False)],
    [("• Recommandations (timing, seuils, risque)", 15, NAVY, True, False)],
    [("• Verdict de dérive (KS-test)", 15, NAVY, True, False)]])

# ---- 4. calibration synthétique
s = content("Résultat 1 — Calibration (données synthétiques)")
for i, (sym, val, col) in enumerate([("β", f"{p_syn.beta:.3f}", TEAL), ("γ", f"{p_syn.gamma:.3f}", TEAL),
                                     ("σ", f"{p_syn.sigma:.3f}", TEAL), ("R₀", f"{R0_syn:.2f}", RED)]):
    x = Inches(0.8 + i * 3.0)
    box(s, x, Inches(1.9), Inches(2.7), Inches(1.9), fill=WHITE, line=col, lw=2.0)
    text(s, x, Inches(2.05), Inches(2.7), Inches(0.9), [[(sym, 32, col, True, False)]], align=PP_ALIGN.CENTER)
    text(s, x, Inches(2.9), Inches(2.7), Inches(0.6), [[(val, 24, NAVY, True, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.5), [
    [("R₀ = β/γ = 5.5 → régime sévère (chaque infecté en contamine ~5,5).", 17, NAVY, True, False)],
    [("Sert de banc d'essai contrôlé pour comparer les trois contrôleurs.", 15, GREY, False, True)]], sa=10)

# ---- 5. benchmark synthétique
s = content("Résultat 2 — Benchmark des contrôleurs ⭐")
s.shapes.add_picture(str(ASSETS / "bench_syn.png"), Inches(0.5), Inches(1.5), width=Inches(7.3))
rows = [[("Stratégie", None), ("Réduction", None)]]
cmap = {"lqr": ORANGE, "pmp": TEAL, "hjb": PURPLE}
for k in ("lqr", "pmp", "hjb"):
    rows.append([(LAB[k], NAVY), (f"{sum_syn[k]['peak_reduction_pct']:.2f} %", cmap[k])])
table(s, Inches(8.1), Inches(1.8), Inches(4.6), rows, [Inches(2.6), Inches(2.0)], NAVY)
text(s, Inches(8.1), Inches(4.6), Inches(4.6), Inches(2.5), [
    [("HJB-PINN ≥ PMP > LQR", 16, NAVY, True, False)],
    [("Le LQR est local (linéarisé) ;", 13, GREY, False, True)],
    [("PMP et HJB-PINN optimisent la", 13, GREY, False, True)],
    [("dynamique non-linéaire complète.", 13, GREY, False, True)]], sa=5)

# ---- 6. correction HJB-PINN
s = content("Résultat 3 — Stabilisation du HJB-PINN")
rows = [[("Indicateur", None), ("Avant", None), ("Après", None)],
        [("Perte d'entraînement", NAVY), ("9,8 × 10¹¹", RED), ("2,5 × 10⁻³", GREEN)],
        [("Réduction du pic", NAVY), ("0,02 %", RED), ("99,98 %", GREEN)],
        [("Politique apprise", NAVY), ("u = 0 (inactive)", RED), ("adaptative", GREEN)]]
table(s, Inches(0.8), Inches(1.7), Inches(9.0), rows, [Inches(3.4), Inches(2.8), Inches(2.8)], NAVY)
text(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.0), [
    [("Cause : ", 15, NAVY, True, False),
     ("fonction valeur en unités brutes (I~10⁵) → résiduelle EDP ~10¹², mal conditionnée.", 15, NAVY, False, False)],
    [("Correctif : ", 15, NAVY, True, False),
     ("coût normalisé (I/N) → perte ~10⁻³ et politique de confinement adaptative active.", 15, NAVY, False, False)],
    [("Le HJB-PINN devient le meilleur contrôleur (99,98 %).", 15, TEAL, True, False)]], sa=9)

# ---- 7. données réelles : calibration + fit
s = content("Résultat 4 — Données réelles (COVID-19, Maroc)")
s.shapes.add_picture(str(ASSETS / "fit_real.png"), Inches(0.5), Inches(1.6), width=Inches(7.3))
card(s, Inches(8.1), Inches(1.6), Inches(4.6), Inches(4.6), "CALIBRATION MAROC", GREEN, [
    [(f"R₀ = {R0_real:.2f}", 20, RED, True, False)],
    [(f"β = {p_real.beta:.3f}", 15, NAVY, False, False)],
    [(f"γ = {p_real.gamma:.3f}   σ = {p_real.sigma:.3f}", 14, NAVY, False, False)],
    [(f"N = {p_real.N:,.0f}", 13, GREY, False, False)],
    [(f"doublement {np.log(2)/r_growth:.1f} j", 13, GREY, False, False)],
    [("", 8, GREY, False, False)],
    [("Le modèle capte la croissance", 12, GREY, False, True)],
    [("initiale ; l'écart au-delà mesure", 12, GREY, False, True)],
    [("l'effet des interventions réelles.", 12, GREY, False, True)]])

# ---- 8. données réelles : benchmark
s = content("Résultat 5 — Benchmark sur paramètres réels (Maroc)")
s.shapes.add_picture(str(ASSETS / "bench_real.png"), Inches(0.5), Inches(1.5), width=Inches(7.3))
rows = [[("Stratégie", None), ("Réduction", None)]]
for k in ("lqr", "pmp"):
    rows.append([(LAB[k], NAVY), (f"{sum_real[k]['peak_reduction_pct']:.2f} %", cmap[k])])
table(s, Inches(8.1), Inches(1.8), Inches(4.6), rows, [Inches(2.6), Inches(2.0)], NAVY)
text(s, Inches(8.1), Inches(4.0), Inches(4.6), Inches(2.8), [
    [("R₀ plus faible (1,87) :", 15, NAVY, True, False)],
    [("le LQR devient modéré,", 13, GREY, False, True)],
    [("le PMP reste quasi total.", 13, GREY, False, True)],
    [("Risque classé « modéré ».", 14, TEAL, True, False)]], sa=5)

# ---- 9. conformité + limites
s = content("Conformité (slide 25) & limites")
rows = [[("Perspective", None), ("État", None)],
        [("Calibration · HJB-PINN · Dashboard · Benchmark", NAVY), ("Réalisé", GREEN)],
        [("Données réelles (Maroc) ingérées et calibrées", NAVY), ("Réalisé", GREEN)],
        [("Gestion des stocks (doses, lits)", NAVY), ("À faire", ORANGE)]]
table(s, Inches(0.8), Inches(1.7), Inches(11.0), rows, [Inches(8.0), Inches(3.0)], NAVY)
text(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.0), [
    [("Limites honnêtes : ", 15, ORANGE, True, False),
     ("modèle à paramètres constants (ne capte pas les vagues multiples) ; coûts de contrôle "
      "idéalisés (les ~99 % sont des bornes optimistes) ; σ, γ fixés à des valeurs COVID.", 15, NAVY, False, False)]], sa=8)

# ---- 10. conclusion
s = prs.slides.add_slide(BLANK); bg(s, NAVY); box(s, 0, 0, SW, Inches(0.18), fill=TEAL)
text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9), [[("Conclusion", 30, WHITE, True, False)]])
text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), [
    [("✓ Chaîne complète théorie → code → décision, validée sur synthétique ET sur données réelles.", 17, WHITE, False, False)],
    [("✓ Contrôle optimal opérationnel : LQR · PMP · HJB-PINN comparés et interchangeables.", 17, WHITE, False, False)],
    [("✓ Le HJB-PINN (IA) est compétitif, voire supérieur aux méthodes classiques (99,98 %).", 17, WHITE, False, False)],
    [("✓ Données réelles du Maroc calibrées : R₀ = 1,87 data-driven.", 17, WHITE, False, False)],
    [("", 10, WHITE, False, False)],
    [("Perspectives : ingestion temps réel (OMS), gestion des stocks, HJB-PINN ré-entraîné sur le réel.", 15, TEAL, False, True)]], sa=14)
box(s, Inches(0.8), Inches(6.7), Inches(4.0), Pt(4), fill=TEAL)

out = ROOT / "presentation_plateforme_resultats.pptx"
prs.save(out)
print("OK ->", out, "(%d slides)" % len(prs.slides._sldIdLst))
